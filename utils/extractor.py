from docx import Document
from langchain_community.document_loaders import PyMuPDFLoader
from pptx import Presentation

from schemas.exception import DocumentExtractionError
from utils.helper import validate_file_extension
from utils.logger import get_logger

logger = get_logger()


class DocumentExtractor:
    """Services for extracting text from various document formats"""

    def __init__(self, file_path: str) -> None:
        self.file_path = file_path
        logger.info("DocumentExtractor initialized")

    def extract(self) -> str:
        file_ext = validate_file_extension(self.file_path)
        try:
            if file_ext == "pdf":
                return self._extract_pdf(self.file_path)
            elif file_ext == "docx":
                return self._extract_docx(self.file_path)
            elif file_ext == "pptx":
                return self._extract_pptx(self.file_path)
            else:
                raise DocumentExtractionError(f"Unsupported file type: {file_ext}")

        except Exception as e:
            logger.error("Extraction failed: ", extra={"file_path": self.file_path})
            raise DocumentExtractionError(f"Failed to extract document: {e}")

    def _extract_pdf(self, file_path: str) -> str:
        try:
            loader = PyMuPDFLoader(file_path, mode="page", extract_tables="markdown")
            documents = loader.load()
            if not documents:
                logger.warning("No content extracted from PDF")
                return ""
            full_text = "\n\n".join(
                [
                    doc.page_content.strip()
                    for doc in documents
                    if doc.page_content.strip()
                ]
            )
            logger.info("PDF extraction successful", extra={"file_path": file_path})
            return full_text
        except Exception as e:
            raise DocumentExtractionError(f"Error extracting PDF: {e}")

    def _extract_docx(self, file_path: str) -> str:
        try:
            doc = Document(file_path)
            paragraphs = [
                para.text.strip() for para in doc.paragraphs if para.text.strip()
            ]
            if not paragraphs:
                logger.warning("No content extracted from DOCX")
                return ""
            full_text = "\n\n".join(paragraphs)
            logger.info("DOCX extraction successful", extra={"file_path": file_path})

            return full_text
        except Exception as e:
            raise DocumentExtractionError(f"Error extracting DOCX: {e}")

    def _extract_pptx(self, file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            text_items = []
            slide_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    text_items.append(f"--- Slide {slide_num} ---")
                    text_items.extend(slide_texts)
                    slide_count += 1

            if not text_items:
                logger.warning(
                    "No text content found in PPTX: ", extra={"file_path": file_path}
                )
                return ""

            full_text = "\n".join(text_items)
            logger.info(
                "Successfully extracted : ",
                extra={"file_path": file_path, "slides_extracted": slide_count},
            )
            return full_text
        except Exception as e:
            raise DocumentExtractionError(f"Error extracting PPTX: {e}")
